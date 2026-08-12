"""Parser prezentacji PowerPoint w formacie Office Open XML (.pptx, .ppsx).

Tresc jest odczytywana slajd po slajdzie, w kolejnosci pokazu. Tytul slajdu
staje sie naglowkiem dziedziczonym przez pozostale sekcje slajdu, a numer
slajdu trafia do pola strony, dzieki czemu wynik wyszukiwania wskazuje
konkretny slajd. Notatki prelegenta sa indeksowane razem ze slajdem.

Parser nie renderuje slajdow, wiec nie zglasza potrzeby OCR: prezentacja
zlozona z samych obrazow jest raportowana jako pusta, z liczba grafik
w komunikacie.
"""

from __future__ import annotations

import datetime as _dt
import re
import zipfile
from collections.abc import Iterator
from pathlib import Path
from typing import Any

from pptx import Presentation as load_presentation
from pptx.exc import PackageNotFoundError
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from finddocs.errors import (
    CorruptedFileError,
    EmptyDocumentError,
    ExtractionError,
    FindDocsError,
    PasswordProtectedError,
)
from finddocs.extractors.base import ExtractionContext, Extractor
from finddocs.normalization.text import clean_text, looks_like_garbage
from finddocs.types import (
    DocumentMetadata,
    ExtractedSection,
    ExtractionResult,
    SupportLevel,
    TextOrigin,
)

#: Sygnatura kontenera OLE. Zaszyfrowany plik .pptx jest zapisany wlasnie tak.
_OLE_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

#: Co ile przetworzonych ksztaltow sprawdzac anulowanie i limit czasu.
_CHECKPOINT_EVERY = 16

#: Maksymalna dlugosc komorki, ktora moze byc jeszcze nazwa kolumny.
_MAX_HEADER_CELL_CHARS = 64

#: Maksymalna dlugosc naglowka przypisywanego sekcjom slajdu.
_MAX_HEADING_CHARS = 200

#: Ile znakow tresci wystarczy do oceny, czy tekst nie jest smieciem.
_GARBAGE_SAMPLE_CHARS = 20_000

#: Wartosc shape_type ksztaltu grupujacego (MSO_SHAPE_TYPE.GROUP).
_SHAPE_TYPE_GROUP = 6

#: Wartosc shape_type obrazu (MSO_SHAPE_TYPE.PICTURE).
_SHAPE_TYPE_PICTURE = 13

_NUMERIC_NOISE_RE = re.compile(r"[\s.,%+\-()]")


class PptxExtractor(Extractor):
    """Adapter prezentacji PowerPoint oparty o biblioteke python-pptx."""

    name = "pptx"
    extensions = (".pptx", ".ppsx", ".pptm", ".ppsm")
    mime_types = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.presentationml.slideshow",
        "application/vnd.ms-powerpoint.presentation.macroEnabled.12",
        "application/vnd.ms-powerpoint.slideshow.macroEnabled.12",
    )
    support_level = SupportLevel.FULL
    priority = 120

    def extract(self, path: Path, context: ExtractionContext) -> ExtractionResult:
        """Wyciaga tekst slajdow, tabele, notatki prelegenta i metadane."""
        context.checkpoint()
        presentation = self._open(path)
        result = ExtractionResult(
            origin=TextOrigin.NATIVE,
            parser_name=self.name,
            support_level=self.support_level,
        )
        try:
            truncated = self._read_slides(presentation, context, result)
        except FindDocsError:
            raise
        except (zipfile.BadZipFile, KeyError, ValueError, IndexError, AttributeError) as exc:
            raise CorruptedFileError(
                "Struktura prezentacji jest uszkodzona, nie udało się odczytać treści.",
                details={"plik": path.name},
            ) from exc
        except OSError as exc:
            raise ExtractionError(
                "Nie udało się odczytać zawartości prezentacji.",
                details={"plik": path.name},
            ) from exc
        except Exception as exc:
            raise ExtractionError(
                "Nieoczekiwany błąd podczas odczytu prezentacji.",
                details={"plik": path.name, "typ_bledu": type(exc).__name__},
            ) from exc

        if not result.sections:
            raise self._empty_error(presentation, path)

        if truncated:
            result.warnings.append(
                "Prezentacja przekroczyła limit długości tekstu, zaindeksowano tylko początek."
            )
        slide_count = len(presentation.slides)
        result.total_pages = slide_count
        result.metadata = self._read_metadata(presentation, result)
        result.metadata.page_count = slide_count
        sample = _sample_text(result.sections)
        if sample and looks_like_garbage(sample):
            result.warnings.append(
                "Tekst prezentacji wygląda na uszkodzony, "
                "udział znaków alfanumerycznych jest niski."
            )
        return result

    # --- otwieranie pliku ------------------------------------------------------

    def _open(self, path: Path) -> Presentation:
        """Otwiera pakiet prezentacji i tlumaczy bledy biblioteki na wyjatki FindDocs."""
        if _is_ole_container(path):
            raise PasswordProtectedError(
                "Prezentacja jest zaszyfrowana lub zabezpieczona hasłem.",
                details={"plik": path.name},
            )
        try:
            return load_presentation(str(path))
        except PackageNotFoundError as exc:
            raise CorruptedFileError(
                "Plik nie jest poprawnym pakietem prezentacji albo jest uszkodzony.",
                details={"plik": path.name},
            ) from exc
        except (zipfile.BadZipFile, KeyError) as exc:
            raise CorruptedFileError(
                "Archiwum prezentacji jest niekompletne, brakuje wymaganych części.",
                details={"plik": path.name},
            ) from exc
        except (AttributeError, IndexError, ValueError) as exc:
            raise CorruptedFileError(
                "Struktura pakietu prezentacji jest nieprawidłowa.",
                details={"plik": path.name},
            ) from exc
        except OSError as exc:
            raise ExtractionError(
                "Nie udało się odczytać pliku prezentacji.",
                details={"plik": path.name},
            ) from exc
        except Exception as exc:
            raise ExtractionError(
                "Nieoczekiwany błąd podczas otwierania prezentacji.",
                details={"plik": path.name, "typ_bledu": type(exc).__name__},
            ) from exc

    def _empty_error(self, presentation: Presentation, path: Path) -> EmptyDocumentError:
        """Buduje wyjatek dla prezentacji bez tresci tekstowej."""
        images = _count_images(presentation)
        if images:
            message = (
                f"Prezentacja nie zawiera tekstu, znaleziono {images} obiektow graficznych. "
                "Tresc jest prawdopodobnie wstawiona jako obraz."
            )
        else:
            message = "Prezentacja nie zawiera żadnej treści tekstowej."
        return EmptyDocumentError(
            message,
            details={"plik": path.name, "obrazy": images, "parser": self.name},
        )

    # --- tresc -----------------------------------------------------------------

    def _read_slides(
        self,
        presentation: Presentation,
        context: ExtractionContext,
        result: ExtractionResult,
    ) -> bool:
        """Przechodzi po slajdach w kolejnosci pokazu. Zwraca, czy tresc obcieto."""
        collector = _Collector(context=context, result=result)
        for number, slide in enumerate(presentation.slides, start=1):
            collector.start_slide(number, _slide_title(slide))
            for shape in _iter_shapes(slide.shapes):
                collector.tick()
                self._add_shape(shape, collector)
                if collector.truncated:
                    return True
            self._add_notes(slide, collector)
            if collector.truncated:
                return True
        return False

    def _add_shape(self, shape: BaseShape, collector: _Collector) -> None:
        """Dodaje tresc jednego ksztaltu: pole tekstowe albo tabele."""
        if getattr(shape, "has_table", False):
            self._add_table(shape, collector)
            return
        if not getattr(shape, "has_text_frame", False):
            return
        for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
            text = "".join(run.text for run in paragraph.runs) or paragraph.text
            collector.add(text)

    def _add_table(self, shape: BaseShape, collector: _Collector) -> None:
        """Zamienia tabele slajdu na sekcje wierszy, jak w dokumentach Word."""
        rows = _table_rows(shape, collector)
        if not rows:
            return
        header = rows[0] if _looks_like_header(rows) else None
        if header is not None:
            collector.add(
                " | ".join(cell for cell in header if cell),
                kind="table_header",
                row=1,
            )
        data_rows = rows[1:] if header is not None else rows
        first_number = 2 if header is not None else 1
        for offset, cells in enumerate(data_rows):
            collector.tick()
            if collector.truncated:
                return
            collector.add(
                _format_row(header, cells),
                kind="table_row",
                row=first_number + offset,
            )

    def _add_notes(self, slide: Slide, collector: _Collector) -> None:
        """Dopisuje notatki prelegenta do sekcji slajdu."""
        try:
            if not slide.has_notes_slide:
                return
            frame = slide.notes_slide.notes_text_frame
        except (KeyError, ValueError, AttributeError):
            return
        if frame is None:
            return
        text = frame.text
        if text.strip():
            collector.add(text, extra={"zrodlo": "notatki prelegenta"})

    # --- metadane --------------------------------------------------------------

    def _read_metadata(
        self, presentation: Presentation, result: ExtractionResult
    ) -> DocumentMetadata:
        """Czyta wlasciwosci pakietu. Braki nie sa bledem, tylko ostrzezeniem."""
        metadata = DocumentMetadata()
        try:
            properties = presentation.core_properties
            metadata.title = _text_or_none(properties.title)
            metadata.author = _text_or_none(properties.author)
            metadata.subject = _text_or_none(properties.subject)
            metadata.keywords = _text_or_none(properties.keywords)
            metadata.language = _text_or_none(properties.language)
            metadata.created_at = _as_datetime(properties.created)
            metadata.modified_at = _as_datetime(properties.modified)
            editor = _text_or_none(properties.last_modified_by)
            if editor:
                metadata.extra["ostatnio_zapisal"] = editor
            category = _text_or_none(properties.category)
            if category:
                metadata.extra["kategoria"] = category
        except (KeyError, ValueError, AttributeError):
            result.warnings.append("Nie udało się odczytać właściwości prezentacji.")
        return metadata


class _Collector:
    """Zbiera sekcje prezentacji z zachowaniem limitow kontekstu ekstrakcji."""

    __slots__ = ("chars", "context", "heading", "page", "result", "steps", "truncated")

    def __init__(self, *, context: ExtractionContext, result: ExtractionResult) -> None:
        self.context = context
        self.result = result
        self.page = 1
        self.heading: str | None = None
        self.chars = 0
        self.truncated = False
        self.steps = 0

    def start_slide(self, number: int, title: str | None) -> None:
        """Ustawia biezacy slajd; jego tytul dziedzicza pozostale sekcje."""
        self.context.checkpoint()
        self.page = number
        self.heading = title[:_MAX_HEADING_CHARS] if title else None

    def tick(self) -> None:
        """Okresowo sprawdza anulowanie i limit czasu."""
        self.steps += 1
        if self.steps % _CHECKPOINT_EVERY == 0:
            self.context.checkpoint()

    def add(
        self,
        text: str,
        *,
        kind: str = "text",
        row: int | None = None,
        extra: dict[str, Any] | None = None,
    ) -> None:
        """Dodaje sekcje po oczyszczeniu tekstu. Puste sekcje sa pomijane."""
        if self.truncated:
            return
        cleaned = clean_text(text)
        if not cleaned:
            return
        remaining = self.context.max_chars - self.chars
        if remaining <= 0:
            self.truncated = True
            return
        if len(cleaned) > remaining:
            cleaned = cleaned[:remaining].rstrip()
            self.truncated = True
            if not cleaned:
                return
        self.chars += len(cleaned)
        self.result.sections.append(
            ExtractedSection(
                text=cleaned,
                kind=kind,
                order=len(self.result.sections),
                page=self.page,
                row=row,
                origin=TextOrigin.NATIVE,
                heading=self.heading,
                extra=extra or {},
            )
        )


# --- przechodzenie po strukturze prezentacji -------------------------------------


def _iter_shapes(shapes: Any) -> Iterator[BaseShape]:
    """Zwraca ksztalty slajdu, wchodzac rekurencyjnie do ksztaltow grupujacych."""
    for shape in shapes:
        if int(getattr(shape, "shape_type", 0) or 0) == _SHAPE_TYPE_GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _slide_title(slide: Slide) -> str | None:
    """Tekst symbolu zastepczego tytulu slajdu, gdy slajd go ma."""
    try:
        title = slide.shapes.title
    except (KeyError, ValueError, AttributeError):
        return None
    if title is None:
        return None
    cleaned = clean_text(title.text)
    return cleaned or None


def _table_rows(shape: BaseShape, collector: _Collector) -> list[list[str]]:
    """Zwraca tresc tabeli jako liste wierszy. Scalone komorki nie sa powielane."""
    rows: list[list[str]] = []
    try:
        table_rows = list(shape.table.rows)  # type: ignore[attr-defined]
    except (ValueError, IndexError, AttributeError):
        return rows
    for row in table_rows:
        collector.tick()
        if collector.truncated:
            break
        values: list[str] = []
        try:
            cells = list(row.cells)
        except (ValueError, IndexError, AttributeError):
            continue
        for cell in cells:
            if getattr(cell, "is_spanned", False):
                continue
            values.append(_flatten(cell.text))
        rows.append(values)
    return rows


def _looks_like_header(rows: list[list[str]]) -> bool:
    """Ocenia, czy pierwszy wiersz tabeli jest wierszem nazw kolumn."""
    if len(rows) < 2:
        return False
    first = rows[0]
    filled = [cell for cell in first if cell]
    if not filled:
        return False
    if len(filled) * 2 < len(first):
        return False
    if any(len(cell) > _MAX_HEADER_CELL_CHARS for cell in filled):
        return False
    if not any(any(ch.isalpha() for ch in cell) for cell in filled):
        return False
    return not any(_is_numeric(cell) for cell in filled)


def _is_numeric(value: str) -> bool:
    """True dla komorek zawierajacych wylacznie liczbe, kwote albo procent."""
    digits = _NUMERIC_NOISE_RE.sub("", value)
    return bool(digits) and digits.isdigit()


def _format_row(header: list[str] | None, cells: list[str]) -> str:
    """Sklada wiersz tabeli w linie tekstu, z nazwami kolumn gdy sa dostepne."""
    parts: list[str] = []
    for index, value in enumerate(cells):
        if not value:
            continue
        label = header[index] if header is not None and index < len(header) else ""
        parts.append(f"{label}: {value}" if label else value)
    return " | ".join(parts)


def _count_images(presentation: Presentation) -> int:
    """Liczy obrazy na slajdach, uzywane w komunikacie o pustej prezentacji."""
    total = 0
    try:
        for slide in presentation.slides:
            for shape in _iter_shapes(slide.shapes):
                if int(getattr(shape, "shape_type", 0) or 0) == _SHAPE_TYPE_PICTURE:
                    total += 1
    except (KeyError, ValueError, AttributeError):
        return total
    return total


# --- drobne narzedzia -----------------------------------------------------------


def _is_ole_container(path: Path) -> bool:
    """Zaszyfrowana prezentacja OOXML jest zapisana jako kontener OLE, nie jako ZIP."""
    try:
        with path.open("rb") as handle:
            return handle.read(len(_OLE_SIGNATURE)) == _OLE_SIGNATURE
    except OSError:
        return False


def _flatten(text: str) -> str:
    """Skleja tekst w jedna linie, uzywane dla komorek tabeli."""
    return " ".join(text.split())


def _text_or_none(value: str | None) -> str | None:
    """Oczyszczona wartosc metadanej albo None, gdy pole jest puste."""
    if not value:
        return None
    cleaned = _flatten(clean_text(value))
    return cleaned or None


def _as_datetime(value: _dt.datetime | None) -> _dt.datetime | None:
    """Przepuszcza tylko poprawne daty, bo wlasciwosci pakietu bywaja uszkodzone."""
    if isinstance(value, _dt.datetime):
        return value
    return None


def _sample_text(sections: list[ExtractedSection], limit: int = _GARBAGE_SAMPLE_CHARS) -> str:
    """Poczatkowy fragment tresci uzywany do oceny jakosci tekstu."""
    parts: list[str] = []
    total = 0
    for section in sections:
        parts.append(section.text)
        total += len(section.text)
        if total >= limit:
            break
    return " ".join(parts)


__all__ = ["PptxExtractor"]
