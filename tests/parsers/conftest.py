"""Fixture tworzace male pliki testowe dla parserow formatow.

Kazdy plik powstaje w ``tmp_path`` w chwili uzycia fixture, dzieki czemu
w repozytorium nie ma zadnych binariow. Pliki DOCX, XLSX i graficzne buduja
biblioteki uzywane przez aplikacje (python-docx, openpyxl, Pillow), pliki PDF
generatory z ``finddocs.demo.generate``, a formaty binarne bez zapisu
(OLE, BIFF2, MSG, DOC) modul ``parser_data``.
"""

from __future__ import annotations

import datetime as _dt
import io
import random
import struct
from collections.abc import Callable
from pathlib import Path

import pytest
from parser_data import (
    DISCLAIMER,
    POLISH_SAMPLE,
    build_biff2_workbook,
    build_legacy_doc,
    build_ole_container,
    build_rtf_document,
)

from finddocs.extractors.base import ExtractionContext

#: Data uzywana we wszystkich metadanych plikow testowych.
FIXED_MOMENT = _dt.datetime(2024, 3, 15, 10, 20, 30)

#: Typy wlasciwosci MAPI potrzebne przy budowie pliku MSG.
PT_UNICODE = 0x001F
PT_BINARY = 0x0102
PT_SYSTIME = 0x0040

#: Rozmiar naglowka strumienia __properties_version1.0 obiektu najwyzszego poziomu.
_TOP_LEVEL_HEADER_SIZE = 32

#: Poczatek epoki FILETIME.
_FILETIME_EPOCH = _dt.datetime(1601, 1, 1, tzinfo=_dt.UTC)


# --- podstawy ------------------------------------------------------------------


@pytest.fixture
def context() -> ExtractionContext:
    """Kontekst ekstrakcji z limitami wystarczajacymi dla malych plikow."""
    return ExtractionContext(max_bytes=8 * 1024 * 1024, max_chars=200_000)


@pytest.fixture
def docs_dir(tmp_path: Path) -> Path:
    """Katalog roboczy na pliki testowe."""
    target = tmp_path / "pliki"
    target.mkdir(parents=True, exist_ok=True)
    return target


@pytest.fixture
def write_file(docs_dir: Path) -> Callable[[str, bytes], Path]:
    """Zapisuje surowe bajty pod podana nazwa i zwraca sciezke."""

    def _write(name: str, data: bytes) -> Path:
        target = docs_dir / name
        target.write_bytes(data)
        return target

    return _write


# --- PDF -----------------------------------------------------------------------


@pytest.fixture
def make_text_pdf(write_file: Callable[[str, bytes], Path]) -> Callable[..., Path]:
    """Tworzy PDF z warstwa tekstowa i metadanymi."""
    from finddocs.demo.generate import build_text_pdf

    def _make(
        name: str = "umowa.pdf", paragraphs: list[str] | None = None, **kwargs: object
    ) -> Path:
        title = str(kwargs.get("title", "Umowa testowa"))
        content = paragraphs if paragraphs is not None else _default_pdf_paragraphs()
        return write_file(name, build_text_pdf(title, content, created=FIXED_MOMENT))

    return _make


def _default_pdf_paragraphs() -> list[str]:
    """Tresc PDF wystarczajaco dluga, zeby parser nie uznal jej za uboga."""
    return [
        POLISH_SAMPLE,
        f"{DISCLAIMER} Numer rachunku 00 1234 5678 9012 3456 7890 1234.",
        "Klient zobowiazuje sie do splaty kwoty 1 234,56 PLN w terminie 30 dni. " * 3,
        "Zalacznikiem do umowy jest harmonogram splat oraz tabela oplat i prowizji. " * 3,
    ]


@pytest.fixture
def text_pdf(make_text_pdf: Callable[..., Path]) -> Path:
    """Jednostronicowy PDF z warstwa tekstowa."""
    return make_text_pdf()


@pytest.fixture
def multipage_pdf(make_text_pdf: Callable[..., Path]) -> Path:
    """PDF o dwoch stronach, uzywany do sprawdzenia numeracji stron."""
    filler = [f"Wiersz numer {index}: {POLISH_SAMPLE}." for index in range(60)]
    return make_text_pdf("wielostronicowy.pdf", [POLISH_SAMPLE, DISCLAIMER, *filler])


@pytest.fixture
def scan_pdf(write_file: Callable[[str, bytes], Path]) -> Path:
    """PDF bez warstwy tekstowej: strona jest jednym obrazem JPEG."""
    from finddocs.demo.generate import build_image_pdf, render_scan_image

    image = render_scan_image([POLISH_SAMPLE, DISCLAIMER], width=420, height=594)
    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", quality=80)
    payload = build_image_pdf(buffer.getvalue(), width=image.width, height=image.height)
    return write_file("skan.pdf", payload)


@pytest.fixture
def broken_pdf(write_file: Callable[[str, bytes], Path]) -> Path:
    """Plik z sygnatura PDF, ale bez poprawnej struktury obiektow."""
    return write_file("uszkodzony.pdf", b"%PDF-1.7\n" + b"\x00\xff" * 400)


# --- DOCX ----------------------------------------------------------------------


@pytest.fixture
def make_docx(docs_dir: Path) -> Callable[..., Path]:
    """Tworzy dokument Word z naglowkami, akapitami i opcjonalna tabela."""
    from docx import Document

    def _make(
        name: str = "notatka.docx",
        *,
        heading: str | None = "Procedura przelewów",
        paragraphs: list[str] | None = None,
        table: tuple[list[str], list[list[str]]] | None = None,
        metadata: bool = True,
    ) -> Path:
        document = Document()
        if metadata:
            properties = document.core_properties
            properties.title = "Procedura przelewów krajowych"
            properties.author = "Łucja Żółw"
            properties.subject = "Obsługa klienta"
            properties.keywords = "przelew; procedura; ćwiczenie"
            properties.language = "pl-PL"
            properties.created = FIXED_MOMENT
            properties.modified = FIXED_MOMENT
        if heading is not None:
            document.add_heading(heading, level=1)
        for text in paragraphs or []:
            document.add_paragraph(text)
        if table is not None:
            header, rows = table
            grid = document.add_table(rows=1, cols=len(header))
            for column, label in enumerate(header):
                grid.rows[0].cells[column].text = label
            for row in rows:
                cells = grid.add_row().cells
                for column, value in enumerate(row):
                    cells[column].text = value
        target = docs_dir / name
        document.save(str(target))
        return target

    return _make


@pytest.fixture
def sample_docx(make_docx: Callable[..., Path]) -> Path:
    """Dokument Word z naglowkiem, akapitami i tabela w tej kolejnosci."""
    return make_docx(
        paragraphs=[POLISH_SAMPLE, DISCLAIMER],
        table=(
            ["Kolumna", "Kwota", "Waluta"],
            [["Wpłata gotówkowa", "1 234,56", "PLN"], ["Przelew wychodzący", "99,00", "PLN"]],
        ),
    )


# --- PPTX ----------------------------------------------------------------------


@pytest.fixture
def make_pptx(docs_dir: Path) -> Callable[..., Path]:
    """Tworzy prezentacje PowerPoint ze slajdow tytul plus tresc."""
    from pptx import Presentation
    from pptx.util import Inches

    def _make(
        name: str = "prezentacja.pptx",
        *,
        slides: list[tuple[str | None, list[str]]] | None = None,
        table: tuple[list[str], list[list[str]]] | None = None,
        notes: list[str] | None = None,
        metadata: bool = True,
    ) -> Path:
        presentation = Presentation()
        title_and_content = presentation.slide_layouts[1]
        for position, (title, paragraphs) in enumerate(slides or []):
            slide = presentation.slides.add_slide(title_and_content)
            if title is not None and slide.shapes.title is not None:
                slide.shapes.title.text = title
            body = slide.placeholders[1].text_frame
            for index, text in enumerate(paragraphs):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = text
            if notes is not None and position < len(notes):
                slide.notes_slide.notes_text_frame.text = notes[position]
        if table is not None:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            header, rows = table
            shape = slide.shapes.add_table(
                len(rows) + 1, len(header), Inches(1), Inches(1), Inches(8), Inches(3)
            )
            for column, label in enumerate(header):
                shape.table.cell(0, column).text = label
            for row_index, row in enumerate(rows, start=1):
                for column, value in enumerate(row):
                    shape.table.cell(row_index, column).text = value
        if metadata:
            properties = presentation.core_properties
            properties.title = "Szkolenie z przelewów"
            properties.author = "Łucja Żółw"
            properties.subject = "Obsługa klienta"
            properties.keywords = "przelew; szkolenie"
            properties.created = FIXED_MOMENT
            properties.modified = FIXED_MOMENT
        target = docs_dir / name
        presentation.save(str(target))
        return target

    return _make


# --- XLSX ----------------------------------------------------------------------


@pytest.fixture
def make_xlsx(docs_dir: Path) -> Callable[..., Path]:
    """Tworzy skoroszyt Excel z podanych arkuszy."""
    from openpyxl import Workbook

    def _make(
        name: str = "zestawienie.xlsx",
        sheets: list[tuple[str, list[list[object]]]] | None = None,
        *,
        metadata: bool = True,
    ) -> Path:
        book = Workbook()
        default = book.active
        for index, (sheet_name, rows) in enumerate(sheets or []):
            sheet = default if index == 0 and default is not None else book.create_sheet()
            sheet.title = sheet_name
            for row in rows:
                sheet.append(row)
        if metadata:
            book.properties.title = "Zestawienie transakcji"
            book.properties.creator = "Łucja Żółw"
            book.properties.subject = "Rozliczenia"
        target = docs_dir / name
        book.save(str(target))
        return target

    return _make


@pytest.fixture
def sample_xlsx(make_xlsx: Callable[..., Path]) -> Path:
    """Skoroszyt z dwoma arkuszami, naglowkiem kolumn, data i liczba calkowita."""
    return make_xlsx(
        sheets=[
            (
                "Transakcje",
                [
                    ["Opis", "Kwota", "Data", "Sztuki"],
                    ["Wpłata gotówkowa", 1234.56, _dt.date(2024, 3, 15), 7],
                    [POLISH_SAMPLE, 99.0, _dt.datetime(2024, 3, 16, 8, 30), 12],
                ],
            ),
            ("Podsumowanie", [["Razem", "Waluta"], ["1 333,56", "PLN"], [DISCLAIMER, ""]]),
        ]
    )


# --- XLS (BIFF2) ---------------------------------------------------------------


@pytest.fixture
def make_xls(write_file: Callable[[str, bytes], Path]) -> Callable[..., Path]:
    """Tworzy plik .xls w formacie BIFF2 czytelny dla xlrd."""

    def _make(name: str = "stary.xls", rows: list[list[object]] | None = None) -> Path:
        return write_file(name, build_biff2_workbook(rows or []))

    return _make


# --- CSV i pliki tekstowe ------------------------------------------------------


@pytest.fixture
def make_csv(docs_dir: Path) -> Callable[..., Path]:
    """Zapisuje plik CSV o podanym separatorze i kodowaniu."""

    def _make(
        name: str,
        rows: list[list[str]],
        *,
        delimiter: str = ";",
        encoding: str = "utf-8",
        newline: str = "\r\n",
    ) -> Path:
        text = newline.join(delimiter.join(row) for row in rows) + newline
        target = docs_dir / name
        target.write_bytes(text.encode(encoding))
        return target

    return _make


@pytest.fixture
def make_text(docs_dir: Path) -> Callable[..., Path]:
    """Zapisuje plik tekstowy w podanym kodowaniu."""

    def _make(name: str, text: str, *, encoding: str = "utf-8") -> Path:
        target = docs_dir / name
        target.write_bytes(text.encode(encoding))
        return target

    return _make


# --- RTF -----------------------------------------------------------------------


@pytest.fixture
def make_rtf(write_file: Callable[[str, bytes], Path]) -> Callable[..., Path]:
    """Tworzy dokument RTF w stronie kodowej cp1250."""

    def _make(name: str = "pismo.rtf", **kwargs: object) -> Path:
        paragraphs = list(kwargs.pop("paragraphs", [POLISH_SAMPLE, DISCLAIMER]))  # type: ignore[arg-type]
        return write_file(name, build_rtf_document(paragraphs, **kwargs))  # type: ignore[arg-type]

    return _make


# --- wiadomosci e-mail ---------------------------------------------------------


@pytest.fixture
def make_eml(docs_dir: Path) -> Callable[..., Path]:
    """Tworzy wiadomosc EML z trescia tekstowa, HTML i zalacznikami."""
    from email.message import EmailMessage
    from email.utils import format_datetime

    def _make(
        name: str = "wiadomosc.eml",
        *,
        subject: str = "Umowa: Zażółć gęślą jaźń",
        sender: str = "Łucja Żółw <lucja@example.test>",
        recipient: str = "Michał Wąsik <michal@example.test>",
        plain: str | None = None,
        html: str | None = None,
        attachments: list[tuple[str, str, bytes]] | None = None,
        nested: bool = False,
    ) -> Path:
        message = EmailMessage()
        message["From"] = sender
        message["To"] = recipient
        message["Subject"] = subject
        message["Date"] = format_datetime(FIXED_MOMENT.replace(tzinfo=_dt.UTC))
        message["Message-ID"] = "<test-0001@example.test>"
        if plain is not None:
            message.set_content(plain)
        if html is not None:
            if plain is None:
                message.set_content("")
            message.add_alternative(html, subtype="html")
        for attachment_name, mime, payload in attachments or []:
            maintype, _, subtype = mime.partition("/")
            message.add_attachment(
                payload, maintype=maintype, subtype=subtype, filename=attachment_name
            )
        if nested:
            inner = EmailMessage()
            inner["From"] = "Grażyna Ćwik <grazyna@example.test>"
            inner["To"] = sender
            inner["Subject"] = "Przekazana wiadomość: ćwiczenie"
            inner.set_content(f"Treść przekazana dalej. {DISCLAIMER}\n")
            message.add_attachment(inner)
        target = docs_dir / name
        target.write_bytes(message.as_bytes())
        return target

    return _make


def filetime(moment: _dt.datetime) -> bytes:
    """Zamienia date na 8-bajtowa wartosc FILETIME uzywana przez MAPI."""
    delta = moment.replace(tzinfo=_dt.UTC) - _FILETIME_EPOCH
    return struct.pack("<Q", int(delta.total_seconds()) * 10_000_000)


def msg_fixed_properties(values: list[tuple[int, int, bytes]]) -> bytes:
    """Sklada strumien ``__properties_version1.0`` obiektu najwyzszego poziomu."""
    out = bytearray(b"\x00" * _TOP_LEVEL_HEADER_SIZE)
    for prop_id, prop_type, payload in values:
        out += struct.pack("<HHI", prop_type, prop_id, 0x06)
        out += payload.ljust(8, b"\x00")[:8]
    return bytes(out)


def msg_stream_name(prop_id: int, prop_type: int) -> str:
    """Nazwa strumienia wlasciwosci MAPI, np. ``__substg1.0_0037001F``."""
    return f"__substg1.0_{prop_id:04X}{prop_type:04X}"


@pytest.fixture
def make_msg(write_file: Callable[[str, bytes], Path]) -> Callable[..., Path]:
    """Tworzy plik MSG jako kontener OLE ze strumieniami wlasciwosci MAPI."""

    def _make(
        name: str = "wiadomosc.msg",
        *,
        properties: dict[tuple[int, int], bytes] | None = None,
        fixed: list[tuple[int, int, bytes]] | None = None,
        attachments: list[dict[tuple[int, int], bytes]] | None = None,
    ) -> Path:
        entries: list[tuple[tuple[str, ...], bytes]] = []
        for (prop_id, prop_type), payload in (properties or {}).items():
            entries.append(((msg_stream_name(prop_id, prop_type),), payload))
        entries.append((("__properties_version1.0",), msg_fixed_properties(fixed or [])))
        for index, attachment in enumerate(attachments or []):
            storage = f"__attach_version1.0_#{index:08X}"
            for (prop_id, prop_type), payload in attachment.items():
                entries.append(((storage, msg_stream_name(prop_id, prop_type)), payload))
        return write_file(name, build_ole_container(entries))

    return _make


# --- obrazy --------------------------------------------------------------------


@pytest.fixture
def make_image(docs_dir: Path) -> Callable[..., Path]:
    """Tworzy maly obraz rastrowy w podanym formacie."""
    from PIL import Image, ImageDraw

    def _make(
        name: str,
        *,
        image_format: str = "PNG",
        size: tuple[int, int] = (160, 90),
        frames: int = 1,
        exif: dict[int, str] | None = None,
    ) -> Path:
        pages = []
        for index in range(max(1, frames)):
            page = Image.new("RGB", size, (250, 248, 244))
            draw = ImageDraw.Draw(page)
            draw.rectangle([4, 4, size[0] - 5, size[1] - 5], outline=(40, 40, 40))
            draw.text((10, 10), f"{POLISH_SAMPLE} {index + 1}", fill=(10, 10, 10))
            pages.append(page)
        target = docs_dir / name
        options: dict[str, object] = {}
        if exif:
            block = Image.Exif()
            for tag, value in exif.items():
                block[tag] = value
            options["exif"] = block
        if frames > 1:
            pages[0].save(
                str(target), format=image_format, save_all=True, append_images=pages[1:], **options
            )
        else:
            pages[0].save(str(target), format=image_format, **options)
        return target

    return _make


# --- archiwa i pliki problematyczne --------------------------------------------


@pytest.fixture
def make_zip(docs_dir: Path) -> Callable[..., Path]:
    """Tworzy zwykle archiwum ZIP z podanych wpisow."""
    import zipfile

    def _make(name: str = "archiwum.zip", entries: list[tuple[str, bytes]] | None = None) -> Path:
        target = docs_dir / name
        with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as archive:
            for entry_name, payload in entries or [("dokument.txt", b"Tresc dokumentu.")]:
                archive.writestr(entry_name, payload)
        return target

    return _make


@pytest.fixture
def make_protected_zip(docs_dir: Path) -> Callable[..., Path]:
    """Tworzy archiwum ZIP z klasycznym szyfrowaniem kazdego wpisu."""
    from finddocs.demo.generate import write_protected_zip

    def _make(
        name: str = "zabezpieczony.zip", entries: list[tuple[str, bytes]] | None = None
    ) -> Path:
        target = docs_dir / name
        write_protected_zip(
            target,
            entries or [("raport.txt", b"Tajne dane testowe.")],
            phrase="tajne-haslo",
            rng=random.Random(1234),
            moment=FIXED_MOMENT,
        )
        return target

    return _make


# --- stary format Word ---------------------------------------------------------


@pytest.fixture
def make_legacy_doc(write_file: Callable[[str, bytes], Path]) -> Callable[..., Path]:
    """Tworzy plik .doc: kontener OLE ze strumieniami WordDocument i 1Table."""

    def _make(
        name: str = "pismo.doc",
        pieces: list[tuple[str, bool]] | None = None,
        *,
        encrypted: bool = False,
    ) -> Path:
        content = pieces if pieces is not None else _default_doc_pieces()
        return write_file(name, build_legacy_doc(content, encrypted=encrypted))

    return _make


def _default_doc_pieces() -> list[tuple[str, bool]]:
    """Fragmenty tekstu: jeden jednobajtowy w cp1252, jeden w UTF-16LE."""
    return [
        ("Pismo okolne numer 00-99\r", True),
        (f"{POLISH_SAMPLE}, ćma i żuraw.\r", False),
        (f"{DISCLAIMER}\r", True),
    ]
